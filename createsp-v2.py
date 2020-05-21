import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
import pptx.util
from pptx.util import Inches, Pt
from pptx_sp import add_slide, create_table, create_table2, iter_cells, resize_table_font

print("NEW COMMAND XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX")
'''
this ver can only process if number of successors is less than 6 and needs slicing of the column into 2
'''

sp_input_path = r'C:\Users\mnajib.suib\Documents\Python\project-2\database\input\sp-ptc1-2020.xlsx'
sp_inputfile = pd.ExcelFile(sp_input_path)
sp_numberofsheet = len(sp_inputfile.sheet_names)
print("Number of sheets: %d" % sp_numberofsheet)

df_list = []
df_idlist = []
# step: put data in every sheet into list of dataframe
for s in range(sp_numberofsheet):
    df_spexcel = pd.read_excel(
        sp_input_path, sheet_name=s, skiprows=2, usecols='F:K', nrows=None).iloc[:, [1, 3, 5]]
    # df_spexcel = df_spexcel.iloc[:, [1, 3, 5]]
    df_id = pd.read_excel(sp_input_path, sheet_name=s,
                          skiprows=2, usecols='A:C', nrows=None).fillna('')
    # df_spexcel = df_spexcel.fillna("")
    df_list.append(df_spexcel)  # for each sheet
    df_idlist.append(df_id)  # for each sheet
for ids in df_idlist:
    print(ids)
# sp excel file is created, df_list contain list of dataframe
with pd.ExcelWriter('spexcel_output1.xlsx') as writer:  # pylint: disable=abstract-class-instantiated
    for n, df in enumerate(df_list):
        df.to_excel(writer, 'Sheet%s' % (n+1), index=False)
    writer.save()

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_content_layout = prs.slide_layouts[1]
sl_title, sh_title = add_slide(prs, title_slide_layout, "Title")

##############################################################################################
##############################################################################################
new_dflist = []
for i, df in enumerate(df_list):
    shname = sp_inputfile.sheet_names[i]
    print("\nNEW SHEET: {}".format(shname))
    colname = list(df.columns)  # list of original column name
    # colname = [i.replace(i[-2:],'') for i in colname] # replace ".1" to ""
    colname = [i.replace(i, str(colname.index(i)))
               for i in colname]  # colname gonna contain [0, 1, 2]
    df.columns = colname  # rename the column of df to [0, 1, 2]

    print(df)
    print("Sheet/dataframe shape: {}, rowcount: {}".format(str(df.shape),
                                                           df.shape[0]))  # tuple of (rowlength, collength)

    j = df.count()  # j is a series that contains length of each column
    seriessize = j.size  # size of series. eg: 3
    j.reset_index(drop=True)  # reindex to numeric value
    # print("series size: {} \ncount of each column:".format(seriessize))
    # print(j)
    j_indexmax = int(j.idxmax())  # index for the highest row count, single val
    valueofhighest = j[j_indexmax]  # value of the max row count
    print("column with most row: line {}, index: {}, value: {}".format(
        j_indexmax+1, j_indexmax, valueofhighest))

    if(j[j_indexmax] > 5):
        # j_indexmax is the column index/name that we wanna use here

        # step: separate the dataframe into 3 dataframe(actually a series) of each column, column = series
        # df_of_dfcol = []
        # for i in range(seriessize):
        #     df_in = df[str(i)]
        #     df_of_dfcol.append(df_in)

        # # step: get the df of column with the most row
        # df_extend = df_of_dfcol[j_indexmax]
        # df_addnew1 = df_extend[:3].dropna(how='all').reset_index(drop=True)
        # df_addnew2 = df_extend[3:].dropna(how='all').reset_index(drop=True)
        # s1 = df_addnew1.str.cat(sep=",")
        # s2 = df_addnew2.str.cat(sep=",")
        # # print("s1: \n{} \ns2: \n{}".format(s1, s2))

        # # step: rearrange the split extended column(series) and other columns into sheet dataframe
        # if(j_indexmax == 0):
        #     df = pd.concat([[df_addnew1], df_addnew2,
        #                     df_of_dfcol[1], df_of_dfcol[2]], axis=1)
        # elif(j_indexmax == 1):
        #     df = pd.concat([df_of_dfcol[0], df_addnew1,
        #                     df_addnew2, df_of_dfcol[2]], axis=1)
        # elif(j_indexmax == 2):
        #     df = pd.concat([df_of_dfcol[0], df_of_dfcol[1],
        #                     df_addnew1, df_addnew2], axis=1)

        # step: drop nan value
        df = df.dropna(how='all')
        print("MSG: df has changed:")
        df = df.fillna('')
        print(df)

        # str_sheetname = "Sheet" + str(i)
        # df.to_excel(w, sheet_name = str_sheetname, index = False)
        # w.save()

        # for column in df:
        #     colcontent = df[column] # content of the column
        #     print(type(colcontent))
        #     li_st = [column] # column = name of column
        #     li_st = li_st + colcontent.values.tolist() # concat the name and the content
    else:
        df = df.fillna('')
        print("MSG: df has no changes, no need to reprint")
        # print(df)

    # step: turn the dataframe into list
    sheet_list2 = []
    for i, column in enumerate(df):
        cc = df[column]
        li_st = [column]
        li_st = li_st + cc.values.tolist()  # list of data in df
        sheet_list2.append(li_st)
    new_dflist.append(df)
    print("list of df sheet:")
    print(sheet_list2)

    # step: turn the id df into list
    # id_list = []
    # for column in df_idlist[i]:
    #     col = df_idlist[i][column]
    #     idl = [column]
    #     idl = idl + col.values.tolist()
    #     id_list.append(idl)
    # print("id list:")
    # print(id_list)
    # print("new_dflist length: {}".format(len(new_dflist)))
    # print(new_dflist)

    slide, shape = add_slide(prs, title_content_layout, "Position: .....")
    print(sheet_list2)
    create_table(slide, sheet_list2, 1, 10)

    start = 0.5
    for i in range(4):  # add_picture(image_file, left, top, width=None, height=None)
        pic = slide.shapes.add_picture('image.jpg', pptx.util.Inches(
            start), pptx.util.Inches(1), width=pptx.util.Inches(1), height=pptx.util.Inches(1))
        start += 1.0
    # create_table2(slide, df, 1, 10)

#######################################################################################


with pd.ExcelWriter('spexcel_output2.xlsx') as writer:  # pylint: disable=abstract-class-instantiated
    for n, df in enumerate(new_dflist):
        df.to_excel(writer, 'Sheet%s' % (n+1), index=False)
    writer.save()

# print("\nPART2")
# spexceloutput2 = pd.ExcelFile('spexcel_output2.xlsx')
# spoutput2list = []
# for s in range(len(spexceloutput2.sheet_names)):
#     df_spexcel2 = pd.read_excel(spexceloutput2, sheet_name=s, nrows=None)
#     # df_spexcel2 = df_spexcel.iloc[:,[1,3,5]]
#     # df_spexcel = df_spexcel.fillna("")
#     print(df_spexcel2)
#     spoutput2list.append(df_spexcel2)
# final_list = []
# for df in spoutput2list:
#     li_st2 = []
#     for column in df:
#         colcontent2 = df[column]
#         li_st2 = [column]
#         li_st2 = li_st2 + colcontent2.values.tolist()
#         final_list.append(li_st2)
# print(final_list)


sl_end, sh_end = add_slide(prs, title_slide_layout, "End")

prs.save('createsp-v2-2.pptx')
