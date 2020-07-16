import pandas as pd
import os
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import requests
import urllib.request


def splitBelow9Rows(column, columnlength):
    breakindex = math.ceil(columnlength/2)
    df_a = column[:breakindex]
    df_b = column[breakindex:]
    return df_a, df_b


def splitBelow15Rows(column, columnlength):
    breakindex = math.ceil(columnlength/3)
    df_a = column[:breakindex]
    df_a2 = column[breakindex:]
    df_b, df_c = splitBelow9Rows(df_a2, len(df_a2))
    return df_a, df_b, df_c


def split_eachline(dflines):
    for i, dfline in enumerate(dflines):
        linelen = len(dfline)
        if linelen > 5:
            if linelen <= 9:
                df_a, df_b = splitBelow9Rows(dfline, linelen)
                splitlist = [df_a, df_b]
                dflines[i] = splitlist
            elif linelen <= 15:
                df_a, df_b, df_c = splitBelow15Rows(dfline, linelen)
                splitlist = [df_a, df_b, df_c]
                dflines[i] = splitlist
            else:
                dfpage1 = dfline.iloc[:15]
                dfpage1 = dfline.iloc[15:]
                # cont later
        else:
            dflines[i] = [dfline]
    return dflines


def separate_pages(eachsheet, typ, case):
    case_name = ['Case: (1+2+3)',
                 'Case: (1+2, 3)',
                 'Case: (1, 2+3)',
                 'Case: (1,2,3)']
    print(case_name[case])
    pagelist = []
    eachpagelist = []
    for i, column in enumerate(eachsheet):
        for eachslicedcolumn in column:
            if typ == 'profile':
                eachslicedcolumn.insert(0, str(i+1))
            eachpagelist.append(eachslicedcolumn)
        if case in [1, 2]:
            if case == 1:
                saveatIndex = [1, 2]
            elif case == 2:
                saveatIndex = [0, 2]
            if i in saveatIndex:
                pagelist.append(eachpagelist)
                eachpagelist = []
        if case == 3:
            pagelist.append(eachpagelist)
            eachpagelist = []
    if case == 0:
        pagelist = [eachpagelist]
    return pagelist


def divide_pages(eachsheet, typ):
    li1, li2, li3 = eachsheet
    pagelist = []
    if len(li1) + len(li2) <= 4:
        if len(li1) + len(li2) + len(li3) <= 4:
            case = 0
        else:
            case = 1
    else:
        if len(li2) + len(li3) <= 4:
            case = 2
        else:
            case = 3
    pagelist = separate_pages(eachsheet, typ, case)
    return pagelist


def make_table(slide, eachpagelist, fontsize):
    col = len(eachpagelist)
    row_lenlist = [len(i) for i in eachpagelist]
    row = max(row_lenlist)

    for column in eachpagelist:
        while len(column) < row:
            column.append('')  # add white spaces to nan/empty cells

    top, left, width, height = Inches(1.2), Inches(
        0.15), Inches(12.5), Inches(0.8)

    # print('row: {}, col: {}'.format(row, col))

    table = slide.shapes.add_table(
        row, col+col, left, top, width, height).table
    for columnindex in range(col):
        for rowindex in range(row):
            table.cell(rowindex, 2*columnindex +
                       1).text = eachpagelist[columnindex][rowindex]

    for tablecolumnindex in range(len(table.columns)):
        if tablecolumnindex % 2 == 0:
            cellwidth = Inches(1.0)
        else:
            cellwidth = Inches(12.6/float(col) - 1.0)
        table.columns[tablecolumnindex].width = cellwidth

    for rowindex in range(row):
        for columnindex in range(2*col):
            cell = table.cell(rowindex, columnindex)
            cell.fill.solid()
            if rowindex == 0:  # header
                cell.fill.fore_color.rgb = RGBColor(0, 177, 169)  # turquoise
            else:  # body
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # grey

    for columnindex in range(col):
        i = columnindex*2
        a = table.cell(0, i)
        a.merge(table.cell(0, i+1))
        if(a.text == '1'):
            a.text = '1st Line'
        elif(a.text == '2'):
            a.text = '2nd Line'
        else:
            a.text = '3rd Line'
    resize_table_font(table, fontsize)


def resize_table_font(table, size):
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)
                run.font.name = 'Arial'


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def add_image(shape, idlist, mediafolderpath):
    # if not os.path.isdir(mediafolderpath):
    #     os.makedirs(mediafolderpath)
    for i, col in enumerate(idlist):
        for j, row in enumerate(col):
            staffid = int(row)
            imgfilename = str(staffid) + '.jpg'
            imgfilepath = os.path.join(mediafolderpath, imgfilename)
            if not (os.path.isfile(imgfilepath)):
                imglink = 'https://talentengine.petronas.com/api/resources/staff/' + \
                    str(staffid) + '/picture'
                response = requests.get(imglink)
                if response.status_code == 200:  # if link exist
                    urllib.request.urlretrieve(imglink, imgfilepath)
                    print('downloading', staffid)
                else:
                    imgfilepath = ''
            if not imgfilepath == '':
                shape.add_picture(imgfilepath, Inches(i + 0.25), Inches(
                    j + 1.6), width=Inches(0.787), height=Inches(0.787))


def add_slide(prs, layout, title):

    prs.slide_width = 11887200
    prs.slide_height = 6686550
    slide = prs.slides.add_slide(layout)
    # slide.shapes.title.text = title + "\nCurrent Incumbent: "
    shape = slide.shapes
    if layout == prs.slide_layouts[0]:
        shape.title.text = title
    if layout == prs.slide_layouts[1]:
        shape.title.text = title
        for shap in shape:
            if not shap.has_text_frame:
                continue
            text_frame = shap.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    if layout == prs.slide_layouts[0]:
        return prs
    else:
        return slide, shape


def getFileInfo(filepath):
    file = pd.ExcelFile(filepath)
    filename = Path(filepath).stem
    sheetname_list = file.sheet_names
    return filename, sheetname_list


def get_dataframeList(filepath, sheetname_list):
    profile_dflist = []
    id_dflist = []
    for sheetname in sheetname_list:
        id_df = pd.read_excel(
            filepath, sheet_name=sheetname, usecols='A:C')
        profile_df = pd.read_excel(
            filepath, sheet_name=sheetname, usecols='D:F')
        profile_dflist.append(profile_df)
        id_dflist.append(id_df)
        # print(data_df, '\n', id_df)
    return id_dflist, profile_dflist


def add_slide_sp(id_list, data_list, presentation, sheetname_list, mediafolder):
    for i, (id_df, profile_df) in enumerate(zip(id_list, data_list)):
        print('\nSP', i+1)
        sheet_dict = {'id': id_df, 'profile': profile_df}

        for key, df_value in sheet_dict.items():
            value_list = [df_value.loc[:, colname].dropna(
                how='all').tolist()
                for colname, column in df_value.iteritems()]
            # print('\n', key, 'raw list', value_list)
            value_splitted = split_eachline(value_list)
            value_pagelist = divide_pages(value_splitted, key)
            # print('by pages', key, ':', value_pagelist)
            sheet_dict[key] = value_pagelist

        idpagelist = sheet_dict['id']
        profilepagelist = sheet_dict['profile']
        title_content_layout = presentation.slide_layouts[1]

        for index, (eachpageid, eachpageprofile) in enumerate(zip(idpagelist, profilepagelist)):
            print('Page', index+1)
            slide, shape = add_slide(
                presentation, title_content_layout, sheetname_list[i])
            make_table(slide, eachpageprofile, 10)
            add_image(shape, eachpageid, mediafolder)
    return presentation


def mainProcess(sp_path, mediafolder):
    outputname = 'succession_planning.pptx'
    filename, sheetname_list = getFileInfo(sp_path)
    id_list, profile_list = get_dataframeList(sp_path, sheetname_list)
    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    presentation = add_slide(presentation, title_slide_layout, filename)
    presentation = add_slide_sp(id_list, profile_list, presentation,
                                sheetname_list, mediafolder)
    presentation.save(outputname)
    return presentation, outputname


# sp_path = r'D:\Documents\Python\project-2\database\input\July\new_sp_input_test.xlsx'
# mediafolder = r'D:\Documents\Python\project-2\media'
# outputname = 'julysp.pptx'

# presentation = mainProcess(sp_path, mediafolder, outputname)
