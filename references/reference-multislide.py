'''
1 inch = 914400 EMU
'''

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
# from sp_dummy import list_dummy
'''
prs = Presentation()
t0, t1, t2 = ['First name', 'Second name', 'Last name'], [
    'John', 'Adam', 'Smith'], ['Alan', 'Paul', 'Walker']
t3, t3, t5 = ['Alex', 'Marina', 'Husin'], [
    'Troy', 'Dallas', 'Sivan'], ['Justin', 'Hailey', 'Brock']
t6 = ['Helena', 'Hairul', 'Irene']
t7 = ['Song', 'Ji', 'Hyo']
# list will not print duplicate of the elements inside the list
tlist = [t0, t1, t2]
'''

def add_slide(prs, layout, title):
    prs.slide_width = 11887200
    prs.slide_height = 6686550
    slide = prs.slides.add_slide(layout)
    # slide.shapes.title.text = title + "\nCurrent Incumbent: "
    shape = slide.shapes
    if layout == prs.slide_layouts[0]:
        slide.shapes.title.text = title
    if layout == prs.slide_layouts[1]:
        for shap in shape:
            if not shap.has_text_frame:
                continue
            text_frame = shap.text_frame
            text_frame.text = title
            # text_frame.vertical_anchor = MSO_ANCHOR.TOP
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    return slide, shape


def create_table(slide, data_list, style, fontsize):
    r = 0
    for dt in data_list:
        if len(dt) > r:
            r = len(dt)
    # rows = len(data_list[0])
    c = len(data_list)
    top = Inches(1.5)
    left = right = Inches(0.3)
    width = Inches(12.0)
    height = Inches(0.8)
    if style == 0:
        # horizontally
        rows = c
        cols = r
        table = slide.shapes.add_table(
            rows, cols, left, top, width, height).table
        for r in range(rows):  # 3
            icol = 0
            while icol < cols:
                table.cell(r, icol).text = data_list[r][icol]
                icol += 1
            print("row " + str(r) + " done")
        '''
        for dl in data_list:
            irow = 0
            while irow < cols:
                table.cell(data_list.index(dl), irow).text = dl[irow]
                irow += 1
            print("row " + str(data_list.index(dl)) + " done")
        '''
    elif style == 1:
        # vertically, col by col
        rows = r
        cols = c
        max_row = 5
        if r < max_row:
            print("difference: " + str(c - r))
            # if no extra row, normal
        elif r > max_row:
            print("difference: " + str(r - c))
            # check the
            # if there are extra rows
            index = 0
            rowlengthmax = 0
            for dl in data_list:
                if len(dl) > rowlengthmax:
                    rowlengthmax = len(dl)
                    index = data_list.index(dl)
            # now index contains the index value of the longest list
            extra_list = []

        table = slide.shapes.add_table(
            rows, cols+cols, left, top, width, height).table
        for r in range(cols):
            irow = 0
            while irow < rows:
                ico = r + (r+1)
                table.cell(irow, ico).text = data_list[r][irow]
                irow += 1
            print("col " + str(r) + " done")
        for t in range(len(table.columns)):
            if t % 2 == 0:
                table.columns[t].width = Inches(1.0)
            else:
                table.columns[t].width = Inches(2.2)
        '''
        # fill row by row
        for r in range(rows):
            icol = 0
            while icol < cols:
                table.cell(r, icol).text = data_list[icol][r]
                icol += 1
            print("row " + str(r) + " done")
        '''
    else:
        pass  # or fill in with other value
    resize_table_font(table, fontsize)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def resize_table_font(table, size):
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)


# title_slide_layout = prs.slide_layouts[0]
# title_content_layout = prs.slide_layouts[1]
# slid1, shap1 = add_slide(prs, title_slide_layout, "Succession Plan")
# slid2, shap2 = add_slide(prs, title_content_layout, "Head of bla")
# slid3, shap3 = add_slide(prs, title_content_layout, "Head of bla2")

# create_table(slid2, list_dummy, 0, 11)
# create_table(slid3, list_dummy, 1, 11)


# prs.save('test-multi.pptx')
