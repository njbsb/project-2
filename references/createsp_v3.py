import pandas as pd
import os
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor


def reindex_column(df):
    colname = list(df.columns)
    colname = [i.replace(i, str(colname.index(i))) for i in colname]
    # colname = [i.replace(i[-2:],'') for i in colname] # replace ".1" to ""
    df.columns = colname
    return df


def df_tolist(df):
    sheet = []
    for column in df:
        cc = df[column]
        list_ = [column]  # with header
        list_ += cc.values.tolist()
        sheet.append(list_)
        # print(li_st)
        # print('------------------')
    return sheet


def getFileInfo(filepath):
    file = pd.ExcelFile(filepath)
    filename = Path(filepath).stem
    sheetname_list = file.sheet_names
    return filename, sheetname_list


def splitBelow9Rows(column, columnlength):
    breakindex = math.ceil(columnlength/2)
    df_a = column.iloc[:breakindex].reset_index(drop=True)
    df_b = column.iloc[breakindex:].reset_index(drop=True)
    return df_a, df_b


def splitBelow15Rows(column, columnlength):
    breakindex = math.ceil(columnlength/3)
    df_a = column.iloc[:breakindex]
    df_a2 = column.iloc[breakindex:]
    df_b, df_c = splitBelow9Rows(df_a2, len(df_a2))
    return df_a, df_b, df_c


def add_slide(prs, layout, title):
    prs.slide_width = 11887200
    prs.slide_height = 6686550
    slide = prs.slides.add_slide(layout)
    # slide.shapes.title.text = title + "\nCurrent Incumbent: "
    shape = slide.shapes
    if layout == prs.slide_layouts[0]:
        slide.shapes.title.text = title
    if layout == prs.slide_layouts[1]:
        for shap in shape:
            if not shap.has_text_frame:
                continue
            text_frame = shap.text_frame
            text_frame.text = title
            # text_frame.vertical_anchor = MSO_ANCHOR.TOP
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    return slide, shape


def make_table(slide, eachpagelist, fontsize):
    col = len(eachpagelist)
    row_lenlist = [len(i) for i in eachpagelist]
    row = max(row_lenlist)

    for column in eachpagelist:
        while len(column) < row:
            column.append('')  # add white spaces

    top, left, width, height = Inches(1.2), Inches(
        0.15), Inches(12.5), Inches(0.8)

    print('row: {}, col: {}'.format(row, col))
    table = slide.shapes.add_table(
        row, col+col, left, top, width, height).table
    for cl in range(col):
        for rw in range(row):
            table.cell(rw, 2*cl + 1).text = eachpagelist[cl][rw]
    for t in range(len(table.columns)):
        if t % 2 == 0:
            cellwidth = Inches(1.0)
        else:
            cellwidth = Inches(12.6/float(col) - 1.0)
        table.columns[t].width = cellwidth
    for rw in range(row):
        for cl in range(2*col):
            cell = table.cell(rw, cl)
            cell.fill.solid()
            if rw == 0:
                cell.fill.fore_color.rgb = RGBColor(0, 177, 169)  # turquoise
            else:
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # grey
    for col in range(col):
        i = col*2
        a = table.cell(0, i)
        a.merge(table.cell(0, i+1))
        if(a.text == '1'):
            a.text = '1st Line'
        elif(a.text == '2'):
            a.text = '2nd Line'
        else:
            a.text = '3rd Line'
    resize_table_font(table, fontsize)


def resize_table_font(table, size):
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def page_3in1(dflist):
    print('Case: line 1, 2, 3 in page 1')
    eachpagelist = []
    for i, li in enumerate(dflist):
        for line in li:
            column = line.tolist()
            column.insert(0, str(i+1))
            eachpagelist.append(column)
    # pagelist.append(eachpagelist)
    pagelist = [eachpagelist]
    return pagelist


def page_separate1(dflist):
    print('Case: line 1 in page 1, line 2 and line 3 in page 2')
    pagelist = []
    for i, li in enumerate(dflist):
        if i == 0 or i == 1:
            eachpagelist = []
        for line in li:
            column = line.tolist()
            column.insert(0, str(i + 1))
            eachpagelist.append(column)
        if i != 1:
            pagelist.append(eachpagelist)
    return pagelist


def page_separate3(dflist):
    print('Case: line 1 and line 2 in page 1, line 3 in page 2')
    pagelist = []
    for i, li in enumerate(dflist):
        if i == 0 or i == 3:
            eachpagelist = []
        for line in li:
            column = line.tolist()
            column.insert(0, str(i + 1))
            eachpagelist.append(column)
        if i != 0:
            pagelist.append(eachpagelist)
    return pagelist


def page_separate123(dflist):
    print('Case: line 1, line 2, line 3 in 3 separate pages')
    pagelist = []
    for i, li in enumerate(dflist):
        eachpagelist = []
        for line in li:
            column = line.tolist()
            column.insert(0, str(i+1))
            eachpagelist.append(column)
        pagelist.append(eachpagelist)
    return pagelist


def add_image(shape, idlist):
    print('idlist ', idlist)


directory = os.getcwd()
filepath = r'D:\Documents\Python\project-2\database\input\SP_GHRM.xlsx'
filename, sheetname_list = getFileInfo(filepath)

data_df_list = []
id_df_list = []
for i, sheetname in enumerate(sheetname_list):
    sheetdf = pd.read_excel(filepath, sheet_name=sheetname,
                            skiprows=2, nrows=None)
    data_df = sheetdf.iloc[:, [6, 8, 10]].dropna(how='all')
    id_df = sheetdf.iloc[:, [0, 1, 2]].dropna(how='all')
    data_df_list.append(data_df)
    id_df_list.append(id_df)

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_content_layout = prs.slide_layouts[1]
sl_title, sh_title = add_slide(prs, title_slide_layout, "TITLE")

for i, data_df in enumerate(data_df_list):
    data_df = reindex_column(data_df)
    id_df_list[i] = reindex_column(id_df_list[i])
    # sheet = df_tolist(data_df)
    rowcount, columncount = data_df.shape
    print("\nCASE %d" % i)
    print(data_df)
    collen_srs = data_df.count()  # contains length of each column
    collen_srs.reset_index(drop=True)
    print("columnlengthsrs: \n{}".format(collen_srs))
    idx_max = int(collen_srs.idxmax())  # which line is the longest
    print("idx max = ", idx_max)
    # longestcolumn = collen_srs[idx_max]  # how long is the line

    columnlenlist = []
    for j in collen_srs:
        columnlenlist.append(j)

    if(idx_max == 1):
        d1, d2, d3 = columnlenlist
        df1 = data_df.iloc[:, 0].dropna(how='all')
        df2 = data_df.iloc[:, 1].dropna(how='all')
        df3 = data_df.iloc[:, 2].dropna(how='all')
        priority_dict = {0: df1, 1: df2, 2: df3}
        # print("indexes: {}, {}, {}".format(1, 2, 3))
        for pd in priority_dict:
            dfpd = priority_dict[pd]
            collen = len(dfpd)
            if collen > 5:
                if collen <= 9:
                    dfa, dfb = splitBelow9Rows(dfpd, collen)
                    splitlist = [dfa, dfb]
                    priority_dict[pd] = splitlist
                else:
                    if collen <= 15:
                        dfa, dfb, dfc = splitBelow15Rows(dfpd, collen)
                        splitlist = [dfa, dfb, dfc]
                        priority_dict[pd] = splitlist
                    else:
                        # data is bigger than 15, need to split to 2
                        dfpage1 = dfpd.iloc[:15]
                        dfpage2 = dfpd.iloc[15:]
                        pass
            else:
                priority_dict[pd] = [dfpd]
        # print(priority_dict)

        # repeat the same as idxmax 0 or 2
    else:
        if(columnlenlist[0] == columnlenlist[2]):
            indexfirst, indexsecond, indexthird = 1, 2, 0
        elif(columnlenlist[0] == columnlenlist[1]):
            indexfirst, indexsecond, indexthird = 2, 1, 0
        else:
            if(idx_max == 0):  # 0
                d1, d2 = columnlenlist[1], columnlenlist[2]
            elif(idx_max == 2):
                d1, d2 = columnlenlist[0], columnlenlist[1]
            dofirst = min(d1, d2)  # length of shortest column
            dosecond = max(d1, d2)  # length of mid column
            dothird = columnlenlist[idx_max]

            indexfirst = columnlenlist.index(dofirst)  # index of shortest
            indexsecond = columnlenlist.index(dosecond)  # index of mid
            indexthird = columnlenlist.index(dothird)  # index of longest
            if(indexsecond == indexfirst):
                for i, cl in enumerate(columnlenlist):
                    if not(i == idx_max or i == columnlenlist.index(dofirst)):
                        indexsecond = i

        # print("indexes: {}, {}, {}".format(
        #     indexfirst, indexsecond, indexthird))

        hold_df = data_df.iloc[:, indexthird].dropna(how='all')
        df1 = data_df.iloc[:, indexfirst].dropna(how='all')
        df2 = data_df.iloc[:, indexsecond].dropna(how='all')

        priority_dict = {indexfirst: df1,
                         indexsecond: df2,
                         indexthird: hold_df}
        priority_id = [indexfirst, indexsecond, indexthird]

        for pd in priority_dict:
            # print("pd: {}".format(pd))
            collen = len(priority_dict[pd])
            dfpd = priority_dict[pd]
            if collen > 5:
                # print("panjang")
                # process split the column into 2 rows here
                if collen <= 9:
                    dfa, dfb = splitBelow9Rows(dfpd, collen)
                    splitlist = [dfa, dfb]
                    priority_dict[pd] = splitlist
                else:
                    if collen <= 15:
                        dfa, dfb, dfc = splitBelow15Rows(dfpd, collen)
                        splitlist = [dfa, dfb, dfc]
                        priority_dict[pd] = splitlist
                    else:
                        if(collen <= 30):
                            # rare case
                            dfpage1 = dfpd.iloc[:15]
                            dfpage2 = dfpd.iloc[15:]
                            # print(dfa, "\n", dfb)
                            dfa, dfb, dfc = splitBelow15Rows(
                                dfpage1, len(dfpage1))
                            splitlist = [dfa, dfb, dfc]
                            if(len(dfpage2) <= 9):
                                dfd, dfe = splitBelow9Rows(
                                    dfpage2, len(dfpage2))
                                splitlist.extend((dfd, dfe))
                            else:
                                dfd, dfe, dff = splitBelow15Rows(
                                    dfpage2, len(dfpage2))
                                splitlist.extend((dfd, dfe, dff))
                            priority_dict[pd] = splitlist
                        else:
                            print("Too much")
                            # dfa do the collen <= 15
                            # dfb do the collen <= 9
            else:
                priority_dict[pd] = [dfpd]

    # all series
    li1, li2, li3 = priority_dict[0], priority_dict[1], priority_dict[2]
    df_list = [li1, li2, li3]
    print(len(li1), len(li2), len(li3))
    if idx_max == 0:
        if len(li2) + len(li3) <= 4:
            if len(li2) + len(li3) + len(li1) <= 4:
                pagelist = page_3in1(df_list)
            else:
                pagelist = page_separate1(df_list)
        else:
            if len(li1) + len(li2) <= 4:
                pagelist = page_separate3(df_list)
            else:
                pagelist = page_separate123(df_list)
    elif idx_max == 1:
        if len(li1) + len(li2) <= 4:
            if len(li1) + len(li2) + len(li3) <= 4:
                pagelist = page_3in1(df_list)
            else:
                pagelist = page_separate3(df_list)
        else:
            if len(li2) + len(li3) <= 4:
                pagelist = page_separate1(df_list)
            else:
                pagelist = page_separate123(df_list)
    elif idx_max == 2:
        if len(li1) + len(li2) <= 4:
            if len(li1) + len(li2) + len(li3) <= 4:
                pagelist = page_3in1(df_list)
            else:
                pagelist = page_separate3(df_list)
        else:
            if len(li2) + len(li3) <= 4:
                pagelist = page_separate1(df_list)
            else:
                pagelist = page_separate123(df_list)
    for page in pagelist:
        slide, shape = add_slide(prs, title_content_layout, sheetname_list[i])
        make_table(slide, page, 10)
        add_image(shape, id_df_list[i])

prs.save("first.pptx")
