from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
# from pptx.enum.dml import MSO_THEME_COLOR


def write_excel(dflist, filename, sheetname):
    with pd.ExcelWriter(filename) as writer:  # pylint: disable=abstract-class-instantiated
        for i, df in enumerate(dflist):
            df.to_excel(writer, sheetname[i], index=False)
        writer.save()


def add_slide(prs, layout, title):
    prs.slide_width = 11887200
    prs.slide_height = 6686550
    slide = prs.slides.add_slide(layout)
    # slide.shapes.title.text = title + "\nCurrent Incumbent: "
    shape = slide.shapes
    if layout == prs.slide_layouts[0]:
        slide.shapes.title.text = title
    if layout == prs.slide_layouts[1]:
        for shap in shape:
            if not shap.has_text_frame:
                continue
            text_frame = shap.text_frame
            text_frame.text = title
            # text_frame.vertical_anchor = MSO_ANCHOR.TOP
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    return slide, shape


def reindex_column(df):
    colname = list(df.columns)
    colname = [i.replace(i, str(colname.index(i))) for i in colname]
    # colname = [i.replace(i[-2:],'') for i in colname] # replace ".1" to ""
    df.columns = colname
    return df


def df_tolist(df):
    sheet = []
    for column in df:
        cc = df[column]
        list_ = [column]  # with header
        list_ += cc.values.tolist()
        sheet.append(list_)
        # print(li_st)
        # print('------------------')
    return sheet


def dfid_tolist(df):
    sheet = []
    for column in df:
        cc = df[column]
        li = cc.values.tolist()
        sheet.append(li)
    return sheet


def normalize_df(df, typ):
    sheet_ = []  # final sheet
    if(df.shape[0] > 5):
        sheet = df_tolist(df)
        srs = df.count()  # length equivalent to column count
        srs.reset_index(drop=True)
        idxmax = int(srs.idxmax())
        maxrowc = srs[idxmax]
        print("column with most row: line {}, index: {}, value: {}".format(
            idxmax+1, idxmax, maxrowc))
        li_extend = sheet[idxmax]
        li_slice1 = li_extend[:4]
        li_slice2 = [str(idxmax)] + li_extend[4:]
        if(idxmax == 0):
            sheet = [li_slice1] + [li_slice2] + [sheet[1]] + [sheet[2]]
        elif(idxmax == 1):
            sheet = [sheet[0]] + [li_slice1] + [li_slice2] + [sheet[2]]
        elif(idxmax == 2):
            sheet = [sheet[0]] + [sheet[1]] + [li_slice1] + [li_slice2]
        df_ = pd.DataFrame.from_records(
            sheet).transpose().dropna(how='all').fillna('')
        if (typ == 0):  # data
            for column in df_:
                cc = df_[column]
                li_ = cc.values.tolist()
                sheet_.append(li_)
                # print(li_)
                # print('------------------')
        elif (typ == 1):  # id
            # ncol = list(df_.iloc[0])  # 0, 1, 2, 2
            # df_.columns = ncol
            df_ = df_[1:]
            sheet_ = dfid_tolist(df_)
    else:  # if df doesnt need to be normalized
        if typ == 0:
            df_ = df.fillna('')
            sheet_ = df_tolist(df_)
        else:
            df_ = df.fillna('')
            sheet_ = dfid_tolist(df_)
    ###
    return df_, sheet_


def create_table(slide, data_list, style, fontsize):
    r = 0
    for dt in data_list:
        if len(dt) > r:
            r = len(dt)
    c = len(data_list)
    top = Inches(1.5)
    left = Inches(0.15)
    width = Inches(12.0)
    height = Inches(0.8)
    if style == 0:  # horizontally
        rows, cols = c, r
        table = slide.shapes.add_table(
            rows, cols, left, top, width, height).table
        for r in range(rows):  # 3
            icol = 0
            while icol < cols:
                table.cell(r, icol).text = data_list[r][icol]
                icol += 1
            print("row %d done" % r)
    elif style == 1:  # vertically, col by col, 1st line - 3rd line
        rows, cols = r, c
        print("rows: {}, cols: {}".format(rows, cols))
        # fill in the table with text data from the list
        table = slide.shapes.add_table(
            rows, cols+cols, left, top, width, height).table
        # fill cells with data from list
        for r in range(cols):
            irow = 0
            while irow < rows:
                ico = r + (r+1)
                table.cell(irow, ico).text = data_list[r][irow]
                irow += 1
            print("col %d done" % r)
        # adjust column width
        for t in range(len(table.columns)):
            if t % 2 == 0:
                table.columns[t].width = Inches(1.0)
            else:
                table.columns[t].width = Inches(2.2)
        # fill cells with colors
        for i in range(rows):
            for j in range(cols+cols):
                a = table.cell(i, j)
                a.fill.solid()
                a.fill.fore_color.rgb = RGBColor(242, 242, 242)  # grey
        # merge header cells, fill header color
        for i in range(len(data_list)):
            i = i*2
            a = table.cell(0, i)
            a.merge(table.cell(0, i+1))
            a.fill.solid()
            a.fill.fore_color.rgb = RGBColor(
                0, 177, 169)  # turquoise
            if(a.text == '0'):
                a.text = '1st Line'
            elif(a.text == '1'):
                a.text = '2nd Line'
            elif(a.text == '2'):
                a.text = '3rd Line'
    resize_table_font(table, fontsize)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def resize_table_font(table, size):
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)
