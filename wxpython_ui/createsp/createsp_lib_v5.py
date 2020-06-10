import pandas as pd
import os
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import requests
import urllib.request


def reindex_column(df):
    colname = list(df.columns)
    colname = [i.replace(i, str(colname.index(i))) for i in colname]
    # colname = [i.replace(i[-2:],'') for i in colname] # replace ".1" to ""
    df.columns = colname
    return df


def df_tolist(df):
    sheet = []
    for column in df:
        cc = df[column]
        list_ = [column]  # with header
        list_ += cc.values.tolist()
        sheet.append(list_)
        # print(li_st)
        # print('------------------')
    return sheet


def getFileInfo(filepath):
    file = pd.ExcelFile(filepath)
    filename = Path(filepath).stem
    sheetname_list = file.sheet_names
    return filename, sheetname_list


def splitBelow9Rows(column, columnlength):
    breakindex = math.ceil(columnlength/2)
    df_a = column.iloc[:breakindex].reset_index(drop=True)
    df_b = column.iloc[breakindex:].reset_index(drop=True)
    return df_a, df_b


def splitBelow15Rows(column, columnlength):
    breakindex = math.ceil(columnlength/3)
    df_a = column.iloc[:breakindex]
    df_a2 = column.iloc[breakindex:]
    df_b, df_c = splitBelow9Rows(df_a2, len(df_a2))
    return df_a, df_b, df_c


def add_slide(prs, layout, title):
    prs.slide_width = 11887200
    prs.slide_height = 6686550
    slide = prs.slides.add_slide(layout)
    # slide.shapes.title.text = title + "\nCurrent Incumbent: "
    shape = slide.shapes
    if layout == prs.slide_layouts[0]:
        slide.shapes.title.text = title
    if layout == prs.slide_layouts[1]:
        for shap in shape:
            if not shap.has_text_frame:
                continue
            text_frame = shap.text_frame
            text_frame.text = title
            # text_frame.vertical_anchor = MSO_ANCHOR.TOP
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
    return slide, shape


def make_table(slide, eachpagelist, fontsize):
    col = len(eachpagelist)
    row_lenlist = [len(i) for i in eachpagelist]
    row = max(row_lenlist)

    for column in eachpagelist:
        while len(column) < row:
            column.append('')  # add white spaces

    top, left, width, height = Inches(1.2), Inches(
        0.15), Inches(12.5), Inches(0.8)

    print('row: {}, col: {}'.format(row, col))
    table = slide.shapes.add_table(
        row, col+col, left, top, width, height).table
    for cl in range(col):
        for rw in range(row):
            table.cell(rw, 2*cl + 1).text = eachpagelist[cl][rw]
    for t in range(len(table.columns)):
        if t % 2 == 0:
            cellwidth = Inches(1.0)
        else:
            cellwidth = Inches(12.6/float(col) - 1.0)
        table.columns[t].width = cellwidth
    for rw in range(row):
        for cl in range(2*col):
            cell = table.cell(rw, cl)
            cell.fill.solid()
            if rw == 0:
                cell.fill.fore_color.rgb = RGBColor(0, 177, 169)  # turquoise
            else:
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # grey
    for col in range(col):
        i = col*2
        a = table.cell(0, i)
        a.merge(table.cell(0, i+1))
        if(a.text == '1'):
            a.text = '1st Line'
        elif(a.text == '2'):
            a.text = '2nd Line'
        else:
            a.text = '3rd Line'
    resize_table_font(table, fontsize)


def resize_table_font(table, size):
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)
                run.font.name = 'Arial'


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def page_3in1(dflist, typ):
    print('Case: line 1, 2, 3 in page 1')
    eachpagelist = []
    for i, li in enumerate(dflist):
        for line in li:
            column = line.tolist()
            if typ == 'sp':
                column.insert(0, str(i+1))
            eachpagelist.append(column)
    pagelist = [eachpagelist]
    return pagelist


def page_separate1(dflist, typ):
    print('Case: line 1 in page 1, line 2 and line 3 in page 2')
    pagelist = []
    for i, li in enumerate(dflist):
        if i == 0 or i == 1:
            eachpagelist = []
        for line in li:
            column = line.tolist()
            if typ == 'sp':
                column.insert(0, str(i + 1))
            eachpagelist.append(column)
        if i != 1:
            pagelist.append(eachpagelist)
    return pagelist


def page_separate3(dflist, typ):
    print('Case: line 1 and line 2 in page 1, line 3 in page 2')
    pagelist = []
    for i, li in enumerate(dflist):
        if i == 0 or i == 3:
            eachpagelist = []
        for line in li:
            column = line.tolist()
            if typ == 'sp':
                column.insert(0, str(i + 1))
            eachpagelist.append(column)
        if i != 0:
            pagelist.append(eachpagelist)
    return pagelist


def page_separate123(dflist, typ):
    print('Case: line 1, line 2, line 3 in 3 separate pages')
    pagelist = []
    for i, li in enumerate(dflist):
        eachpagelist = []
        for line in li:
            column = line.tolist()
            if typ == 'sp':
                column.insert(0, str(i+1))
            eachpagelist.append(column)
        pagelist.append(eachpagelist)
    return pagelist


def manage_dfline(dflines):
    for i, dfline in enumerate(dflines):
        linelen = len(dfline)
        if linelen > 5:
            if linelen <= 9:
                df_a, df_b = splitBelow9Rows(dfline, linelen)
                splitlist = [df_a, df_b]
                dflines[i] = splitlist
            elif linelen <= 15:
                df_a, df_b, df_c = splitBelow15Rows(dfline, linelen)
                splitlist = [df_a, df_b, df_c]
                dflines[i] = splitlist
            else:
                dfpage1 = dfline.iloc[:15]
                dfpage1 = dfline.iloc[15:]
                # cont later
        else:
            dflines[i] = [dfline]
    return dflines


def add_image(shape, idlist, mediafolderpath):
    for i, col in enumerate(idlist):
        for j, row in enumerate(col):
            staffid = int(row)
            imgfilename = str(staffid) + '.jpg'
            imgfilepath = os.path.join(mediafolderpath, imgfilename)
            if not (os.path.isfile(imgfilepath)):
                imglink = 'https://talentengine.petronas.com/api/resources/staff/' + \
                    str(staffid) + '/picture'
                response = requests.get(imglink)
                if response.status_code == 200:  # if link exist
                    urllib.request.urlretrieve(imglink, imgfilepath)
                else:
                    imgfilename = 'default.jpg'
            # else:
            #     imgfilename = os.path.join(mediafolderpath, imgfilename)
            pic = shape.add_picture(imgfilepath, Inches(i + 0.25), Inches(
                j + 1.6), width=Inches(0.787), height=Inches(0.787))


def divide_pages(dflines, idx_max, typ):
    li1, li2, li3 = dflines
    pagelist = []
    if idx_max == 0:
        if len(li2) + len(li3) <= 4:
            if len(li2) + len(li3) + len(li1) <= 4:
                pagelist = page_3in1(dflines, typ)
            else:
                pagelist = page_separate1(dflines, typ)
        else:
            if len(li1) + len(li2) <= 4:
                pagelist = page_separate3(dflines, typ)
            else:
                pagelist = page_separate123(dflines, typ)
    elif idx_max == 1:
        if len(li1) + len(li2) <= 4:
            if len(li1) + len(li2) + len(li3) <= 4:
                pagelist = page_3in1(dflines, typ)
            else:
                pagelist = page_separate3(dflines, typ)
        else:
            if len(li2) + len(li3) <= 4:
                pagelist = page_separate1(dflines, typ)
            else:
                pagelist = page_separate123(dflines, typ)
    elif idx_max == 2:
        if len(li1) + len(li2) <= 4:
            if len(li1) + len(li2) + len(li3) <= 4:
                pagelist = page_3in1(dflines, typ)
            else:
                pagelist = page_separate3(dflines, typ)
        else:
            if len(li2) + len(li3) <= 4:
                pagelist = page_separate1(dflines, typ)
            else:
                pagelist = page_separate123(dflines, typ)
    return pagelist


def get_rawDataframe(filepath, sheetname_list):
    data_df_list = []
    id_df_list = []
    for sheetname in sheetname_list:
        sheetdf = pd.read_excel(filepath, sheet_name=sheetname,
                                skiprows=2, nrows=None)
        data_df = sheetdf.iloc[:, [6, 8, 10]].dropna(how='all')
        id_df = sheetdf.iloc[:, [0, 1, 2]].dropna(how='all')
        data_df_list.append(data_df)
        id_df_list.append(id_df)
    return data_df_list, id_df_list


def createPPT_data(data_df_list, id_df_list, prs, sheetname_list, mediafolderpath):
    for i, data_df in enumerate(data_df_list):
        data_df = reindex_column(data_df)
        id_df = reindex_column(id_df_list[i])
        # sheet = df_tolist(data_df)
        rowcount, columncount = data_df.shape

        collen_srs = data_df.count()  # contains length of each column
        collen_srs.reset_index(drop=True)
        idx_max = int(collen_srs.idxmax())  # which line is the longest
        print("\nCASE {}\n".format(i), data_df,
              "\ncolumnlengthsrs: \n{}".format(collen_srs),
              "\nidx max = ", idx_max)

        columnlenlist = []
        for j in collen_srs:
            columnlenlist.append(j)

        dflines = [data_df.loc[:, colname].dropna(
            how='all') for colname, column in data_df.iteritems()]
        idlines = [id_df.loc[:, colname].dropna(
            how='all') for colname, column in id_df.iteritems()]

        dflines = manage_dfline(dflines)
        idlines = manage_dfline(idlines)

        pagelist = divide_pages(dflines, idx_max, 'sp')
        idpagelist = divide_pages(idlines, idx_max, 'id')
        title_slide_layout = prs.slide_layouts[0]
        title_content_layout = prs.slide_layouts[1]
        for pageid, page in zip(idpagelist, pagelist):
            slide, shape = add_slide(
                prs, title_content_layout, sheetname_list[i])
            make_table(slide, page, 10)
            add_image(shape, pageid, mediafolderpath)


# directory = os.getcwd()
# # filepath = r'D:\Documents\Python\project-2\database\input\SP_GHRM.xlsx'
# filepath = r'D:\Documents\Python\project-2\wxpython_ui\createsp\vp_sp_last.xlsx'
# mediafolderpath = r'D:\Documents\Python\project-2\media'
# filename, sheetname_list = getFileInfo(filepath)

# data_df_list, id_df_list = get_rawDataframe(filepath, sheetname_list)

# prs = Presentation()
# title_slide_layout = prs.slide_layouts[0]
# title_content_layout = prs.slide_layouts[1]
# # sl_title, sh_title = add_slide(prs, title_slide_layout, "TITLE")

# createPPT_data(data_df_list, id_df_list)

# prs.save("first_v4.pptx")
